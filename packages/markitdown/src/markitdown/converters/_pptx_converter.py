import sys
import base64
import os
import io
import re
import html

from typing import BinaryIO, Any
from operator import attrgetter

from ._html_converter import HtmlConverter
from ._llm_caption import llm_caption
from .._base_converter import DocumentConverter, DocumentConverterResult
from .._stream_info import StreamInfo
from .._exceptions import MissingDependencyException, MISSING_DEPENDENCY_MESSAGE

# Try loading optional (but in this case, required) dependencies
# Save reporting of any exceptions for later
_dependency_exc_info = None
try:
    import pptx
except ImportError:
    # Preserve the error and stack trace for later
    _dependency_exc_info = sys.exc_info()


ACCEPTED_MIME_TYPE_PREFIXES = [
    "application/vnd.openxmlformats-officedocument.presentationml",
]


def _sniff_content_type(blob: bytes) -> str:
    """Cheap magic-byte sniffing for image formats PIL won't recognize.

    Returns an image/* mimetype string. python-pptx's `shape.image.content_type`
    asks PIL to identify the blob, which fails (UnidentifiedImageError) for
    EMF/WMF and other Windows-only formats. When that happens we fall back
    here.
    """
    if not blob:
        return "application/octet-stream"
    if blob.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if blob[:3] == b"\xff\xd8\xff":
        return "image/jpeg"
    if blob[:6] in (b"GIF87a", b"GIF89a"):
        return "image/gif"
    if blob[:2] == b"BM":
        return "image/bmp"
    if blob[:4] == b"RIFF" and blob[8:12] == b"WEBP":
        return "image/webp"
    if blob[:4] == b"II*\x00" or blob[:4] == b"MM\x00*":
        return "image/tiff"
    # EMF: starts with EMR_HEADER record type 0x00000001, identifier 'EMF '
    # at offset 40.
    if blob[:4] == b"\x01\x00\x00\x00" and blob[40:44] == b" EMF":
        return "image/x-emf"
    # Placeable WMF
    if blob[:4] == b"\xd7\xcd\xc6\x9a":
        return "image/x-wmf"
    # Non-placeable WMF (less reliable: starts with 0x0001 0x0009)
    if blob[:4] == b"\x01\x00\x09\x00":
        return "image/x-wmf"
    return "application/octet-stream"


def _safe_image_blob_and_type(shape):
    """Read a picture shape's blob + content-type, surviving WMF/EMF/etc.

    Returns (blob, content_type) on success, or (None, None) if the image
    cannot be read at all (e.g. linked external file).
    """
    try:
        blob = shape.image.blob
    except Exception:
        return None, None
    try:
        content_type = shape.image.content_type or _sniff_content_type(blob)
    except Exception:
        # PIL_Image.open() inside python-pptx failed (UnidentifiedImageError
        # for EMF/WMF) — fall back to magic-byte sniffing.
        content_type = _sniff_content_type(blob)
    return blob, content_type

ACCEPTED_FILE_EXTENSIONS = [".pptx"]


class PptxConverter(DocumentConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(self):
        super().__init__()
        self._html_converter = HtmlConverter()

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> bool:
        mimetype = (stream_info.mimetype or "").lower()
        extension = (stream_info.extension or "").lower()

        if extension in ACCEPTED_FILE_EXTENSIONS:
            return True

        for prefix in ACCEPTED_MIME_TYPE_PREFIXES:
            if mimetype.startswith(prefix):
                return True

        return False

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> DocumentConverterResult:
        # Check the dependencies
        if _dependency_exc_info is not None:
            raise MissingDependencyException(
                MISSING_DEPENDENCY_MESSAGE.format(
                    converter=type(self).__name__,
                    extension=".pptx",
                    feature="pptx",
                )
            ) from _dependency_exc_info[
                1
            ].with_traceback(  # type: ignore[union-attr]
                _dependency_exc_info[2]
            )

        # Set up image extraction if enabled
        image_extractor = None
        if kwargs.get("extract_images", False):
            images_dir = kwargs.get("images_dir")
            if images_dir:
                from ..converter_utils._image_extractor import ImageExtractor

                rel_prefix = os.path.basename(images_dir)
                name_prefix = kwargs.get("image_name_prefix", "")
                image_extractor = ImageExtractor(
                    images_dir, rel_prefix=rel_prefix, name_prefix=name_prefix
                )

        # Perform the conversion
        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title

            def get_shape_content(shape, **kwargs):
                nonlocal md_content
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069

                    # Read the image blob and content-type once, surviving
                    # WMF/EMF/etc. that python-pptx can't classify via PIL.
                    blob, content_type = _safe_image_blob_and_type(shape)
                    if blob is None:
                        # Picture has no readable backing (e.g. linked file
                        # that's missing). Skip — don't kill the whole deck.
                        return

                    llm_description = ""
                    alt_text = ""

                    # Potentially generate a description using an LLM
                    llm_client = kwargs.get("llm_client")
                    llm_model = kwargs.get("llm_model")
                    if llm_client is not None and llm_model is not None:
                        # Prepare a file_stream and stream_info for the image data
                        try:
                            image_filename = shape.image.filename
                        except Exception:
                            image_filename = None
                        image_extension = None
                        if image_filename:
                            image_extension = os.path.splitext(image_filename)[1]
                        image_stream_info = StreamInfo(
                            mimetype=content_type,
                            extension=image_extension,
                            filename=image_filename,
                        )

                        image_stream = io.BytesIO(blob)

                        # Caption the image
                        try:
                            llm_description = llm_caption(
                                image_stream,
                                image_stream_info,
                                client=llm_client,
                                model=llm_model,
                                prompt=kwargs.get("llm_prompt"),
                            )
                        except Exception:
                            # Unable to generate a description
                            pass

                    # Also grab any description embedded in the deck
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        # Unable to get alt text
                        pass

                    # Prepare the alt, escaping any special characters
                    alt_text = "\n".join([llm_description, alt_text]) or shape.name
                    alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text)
                    alt_text = re.sub(r"\s+", " ", alt_text).strip()

                    # Image handling: extract > data URI > placeholder
                    if image_extractor is not None:
                        base_name = re.sub(r"\W", "_", shape.name) if shape.name else None
                        filename = image_extractor.save_image(
                            blob,
                            content_type=content_type or "image/png",
                            base_name=base_name,
                        )
                        md_content += f"\n![{alt_text}]({filename})\n"
                    elif kwargs.get("keep_data_uris", False):
                        b64_string = base64.b64encode(blob).decode("utf-8")
                        md_content += f"\n![{alt_text}](data:{content_type or 'image/png'};base64,{b64_string})\n"
                    else:
                        # A placeholder name
                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        md_content += "\n![" + alt_text + "](" + filename + ")\n"

                # Tables
                if self._is_table(shape):
                    md_content += self._convert_table_to_markdown(shape.table, **kwargs)

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        # Write the table as HTML, then convert it to Markdown
        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(html_table, **kwargs).markdown.strip()
            + "\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            # Handle the specific error for unsupported chart types
            if "unsupported plot type" in str(e):
                return "\n\n[unsupported chart]\n\n"
        except Exception:
            # Catch any other exceptions that might occur
            return "\n\n[unsupported chart]\n\n"
